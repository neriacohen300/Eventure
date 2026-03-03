import os
import tkinter as tk
from tkinter import filedialog
from pptx import Presentation
from openpyxl import Workbook
from openpyxl.utils import get_column_letter
import re
from datetime import datetime

# Function to clean text for Excel (remove illegal characters)
# Function to clean text for Excel (remove only illegal control characters)
def clean_text_for_excel(text):
    # Remove non-printable control characters (excluding valid ones like space and punctuation)
    illegal_chars = r"[\x00-\x1F\x7F]"  # Control characters (0x00-0x1F, 0x7F)
    cleaned_text = re.sub(illegal_chars, "", text)  # Remove illegal characters
    return cleaned_text

# Function to extract images from slides and save them
def extract_images_from_slide(slide, pptx_folder, slide_num):
    images = []
    image_shapes = []  # Keep track of image shapes for later checks
    for shape in slide.shapes:
        try:
            if hasattr(shape, "image"):  # Check if the shape has an image attribute
                img = shape.image
                img_bytes = img.blob
                img_format = img.ext
                img_name = f"Slide_{slide_num}_image_{len(images)}.{img_format}"
                img_path = os.path.join(pptx_folder, img_name)

                # Save the image to the folder
                with open(img_path, 'wb') as img_file:
                    img_file.write(img_bytes)

                images.append(img_path)
                # Store the image shape for later checking
                image_shapes.append(shape)
        except ValueError:
            # Skip shapes that do not contain an image (e.g., text boxes)
            pass
    return images, image_shapes

# Function to extract text from a slide
def extract_text_from_slide(slide, slide_num, image_shapes):
    text = ""
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip() != "":
            if not hasattr(shape, "image"):  # Make sure it's not an image with text
                text += f"{shape.text}\n"
            else:
                # Check if the text is within the area of an image shape
                if not is_text_near_image(shape, image_shapes):
                    text += f"{shape.text}\n"
    return text

# Function to check if text is near an image
def is_text_near_image(text_shape, image_shapes):
    for img_shape in image_shapes:
        # Get coordinates and size of the text and image
        text_left, text_top, text_width, text_height = text_shape.left, text_shape.top, text_shape.width, text_shape.height
        img_left, img_top, img_width, img_height = img_shape.left, img_shape.top, img_shape.width, img_shape.height

        # Simple overlap check (you can adjust the margin threshold as needed)
        if (text_left < img_left + img_width and text_left + text_width > img_left and
            text_top < img_top + img_height and text_top + text_height > img_top):
            return True  # Text is near or overlapping with the image
    return False

# Function to extract PPTX content and save it to an Excel sheet
def extract_pptx_content_to_excel(pptx_file):
    pptx_file_name = os.path.splitext(os.path.basename(pptx_file))[0]
    pptx_dir = os.path.dirname(pptx_file)  # Get the directory of the PPTX file
    pptx_folder = os.path.join(pptx_dir, pptx_file_name)  # Create a folder in the same location

    # Create the main folder for this presentation
    if not os.path.exists(pptx_folder):
        os.makedirs(pptx_folder)

    # Create a new workbook and add a worksheet
    wb = Workbook()
    ws = wb.active
    ws.title = "PPTX Extracted Content"

    # Add headers to the Excel sheet
    ws.append(["Slide Number", "Text Content", "Image Paths"])

    # Open the .pptx file
    presentation = Presentation(pptx_file)
    
    for slide_num, slide in enumerate(presentation.slides, start=1):
        # Extract images
        images, image_shapes = extract_images_from_slide(slide, pptx_folder, slide_num)

        # Extract text
        text = extract_text_from_slide(slide, slide_num, image_shapes)
        # Clean text before adding it to the Excel sheet
        text = clean_text_for_excel(text)

        # Add data to the Excel sheet
        ws.append([slide_num, text if text else "No text content", ", ".join(images) if images else "No images"])

        print(f"Processed Slide {slide_num}: {pptx_folder} (Images: {len(images)}, Text: {len(text)} characters)")

    # Save the workbook to an Excel file
    excel_filename = os.path.join(pptx_folder, f"{pptx_file_name}_extracted_content.xlsx")
    wb.save(excel_filename)
    print(f"Excel file saved at: {excel_filename}")


def extract_pptx_content_to_slideshow_file(pptx_file):
    pptx_file_name = os.path.splitext(os.path.basename(pptx_file))[0]
    pptx_dir = os.path.dirname(pptx_file)  # Get the directory of the PPTX file
    pptx_folder = os.path.join(pptx_dir, pptx_file_name)  # Create a folder in the same location

    # Create the main folder for this presentation
    if not os.path.exists(pptx_folder):
        os.makedirs(pptx_folder)

    # Open the .pptx file
    presentation = Presentation(pptx_file)

    current_date = datetime.now().strftime('%Y-%m-%d %H:%M:%S')


    # Prepare the output slideshow file
    slideshow_file = os.path.join(pptx_folder, f"{pptx_file_name}.slideshow")
    with open(slideshow_file, "w", encoding="utf-8") as file:
        file.write("0\n")
        for slide_num, slide in enumerate(presentation.slides, start=1):
            # Extract images
            images, image_shapes = extract_images_from_slide(slide, pptx_folder, slide_num)

            # Extract text
            text = extract_text_from_slide(slide, slide_num, image_shapes)
            # Clean text before using
            text = clean_text_for_excel(text) if text else ""
            
            # Write to the slideshow file
            for i, image in enumerate(images):
                # Only the second image (index 1) gets "True"
                is_second_image = "True" if i == 1 else "False"

                file.write(
                    f"{image},5,fade,1,{text if i == 0 else ''},0,{is_second_image},{current_date},\n"
                )

            print(f"Processed Slide {slide_num}: {pptx_folder} (Images: {len(images)}, Text: {len(text)} characters)")

    print(f"Slideshow file saved at: {slideshow_file}")
    return slideshow_file


# Function to choose a .pptx file using a file dialog
def choose_pptx_file():
    root = tk.Tk()
    root.withdraw()  # Hide the main tkinter window
    pptx_file = filedialog.askopenfilename(
        title="Select a PowerPoint file",
        filetypes=[("PowerPoint files", "*.pptx;*.pptm")]
    )
    return pptx_file

# Main function to run the script
def main():
    pptx_file = choose_pptx_file()
    if pptx_file:
        #extract_pptx_content_to_excel(pptx_file)
        extract_pptx_content_to_slideshow_file(pptx_file)
    else:
        print("No file selected. Exiting.")

# Run the main function
if __name__ == "__main__":
    main()